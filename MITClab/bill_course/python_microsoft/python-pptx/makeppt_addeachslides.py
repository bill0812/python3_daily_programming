import os
from pptx import Presentation
from pptx.util import Inches, Pt

img_path = 'monty-truth.png'

prs = Presentation()

for i in range(len(prs.slide_layouts)) :
    slide = prs.slides.add_slide(prs.slide_layouts[i])
    print(len(slide.shapes.placeholders))

# 儲存 PPT
prs.save(os.getcwd() + '/add_each_slides.pptx')
