from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = ChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Q1 Sales', (19.2, 21.4, 16.7))
chart_data.add_series('Q2 Sales', (22.3, 28.6, 15.2))
chart_data.add_series('Q3 Sales', (20.4, 26.3, 14.2))

graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

chart = graphic_frame.chart

# Axes
category_axis = chart.category_axis
category_axis.has_major_gridlines = True
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(24)

value_axis = chart.value_axis
value_axis.maximum_scale = 50.0
value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
value_axis.has_minor_gridlines = True

tick_labels = value_axis.tick_labels
tick_labels.number_format = '0"%"'
tick_labels.font.bold = True
tick_labels.font.size = Pt(14)

# Data Labels
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels

data_labels.font.size = Pt(13)
data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
data_labels.position = XL_LABEL_POSITION.INSIDE_END

# legend
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = False

prs.save('add_chart05.pptx')