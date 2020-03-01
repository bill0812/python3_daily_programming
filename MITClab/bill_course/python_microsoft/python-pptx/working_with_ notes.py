from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

paragraph_strs = [
    'Egg, bacon, sausage and spam.',
    'Spam, bacon, sausage and spam.',
    'Spam, egg, spam, spam, bacon and spam.'
]

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

print(slide.has_notes_slide)

notes_slide = slide.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = 'foobar'

notes_placeholder = notes_slide.notes_placeholder

for placeholder in notes_slide.placeholders:
    print placeholder.placeholder_format.type

for shape in notes_slide.shapes:
    print shape