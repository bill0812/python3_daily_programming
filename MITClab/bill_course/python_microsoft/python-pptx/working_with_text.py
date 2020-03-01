from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

paragraph_strs = [
    'Egg, bacon, sausage and spam.',
    'Spam, bacon, sausage and spam.',
    'Spam, egg, spam, spam, bacon and spam.'
]

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
shape = slide.shapes.add_textbox(left, top, width, height)

text_frame = shape.text_frame
text_frame.clear()  # remove any existing paragraphs, leaving one empty one

p = text_frame.paragraphs[0]
p.text = paragraph_strs[0]

for para_str in paragraph_strs[1:]:
    p = text_frame.add_paragraph()
    p.text = para_str

shape.text = 'foobar'

# is equivalent to ...

text_frame = shape.text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = 'foobar'

text_frame = shape.text_frame
text_frame.text = 'Spam, eggs, and spam'
text_frame.margin_bottom = Inches(0.08)
text_frame.margin_left = 0
text_frame.vertical_anchor = MSO_ANCHOR.TOP
text_frame.word_wrap = False
text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

text_frame = shape.text_frame
text_frame.clear()

p = text_frame.paragraphs[0]
p.text = paragraph_strs[0]
p.alignment = PP_ALIGN.LEFT

for para_str in paragraph_strs[1:]:
    p = text_frame.add_paragraph()
    p.text = para_str
    p.alignment = PP_ALIGN.LEFT
    p.level = 1

text_frame = shape.text_frame
text_frame.clear()  # not necessary for newly-created shape

p = text_frame.paragraphs[0]
run = p.add_run()
run.text = 'Spam, eggs, and spam'

font = run.font
font.name = 'Calibri'
font.size = Pt(18)
font.bold = True
font.italic = None  # cause value to be inherited from theme
font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
# font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
# run.hyperlink.address = 'https://github.com/scanny/python-pptx'