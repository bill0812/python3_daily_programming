import os
from pptx import Presentation
from pptx.util import Inches, Pt

img_path = 'monty-truth.png'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(5)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

# 儲存 PPT
prs.save(os.getcwd() + '/add_picture.pptx')
