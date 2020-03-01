from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

paragraph_strs = [
    'Egg, bacon, sausage and spam.',
    'Spam, bacon, sausage and spam.',
    'Spam, egg, spam, spam, bacon and spam.'
]

# ---create presentation with 1 slide---
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# ---add table to slide---
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
shape = slide.shapes.add_table(3, 3, x, y, cx, cy)
print(shape)
print(shape.has_table)

table = shape.table
print(table)

slide = prs.slides.add_slide(prs.slide_layouts[2])
table_placeholder = slide.shapes[1]
shape = table_placeholder.insert_table(rows=3, cols=4)
table = shape.table

# access cells
cell = table.cell(0, 0)
print(cell.text)

cell.text = 'Unladen Swallow'

cell = table.cell(0, 0)
other_cell = table.cell(1, 1)
print(cell.is_merge_origin)

cell.merge(other_cell)
print(cell.is_merge_origin)
print(cell.is_spanned)
print(other_cell.is_spanned)
print(table.cell(0, 1).is_spanned)

cell = table.cell(0, 0)
cell.split()
print(cell.is_merge_origin)
print(table.cell(0, 1).is_spanned)

def iter_merge_origins(table):
    """Generate each merge-origin cell in *table*.

    Cell objects are ordered by their position in the table,
    left-to-right, top-to-bottom.
    """
    return (cell for cell in table.iter_cells() if cell.is_merge_origin)

def merged_cell_report(cell):
    """Return str summarizing position and size of merged *cell*."""
    return (
        'merged cell at row %d, col %d, %d cells high and %d cells wide'
        % (cell.row_idx, cell.col_idx, cell.span_height, cell.span_width)
    )

# ---Print a summary line for each merged cell in *table*.---
for merge_origin_cell in iter_merge_origins(table):
    print(merged_cell_report(merge_origin_cell))

def iter_visible_cells(table):
    return (cell for cell in table.iter_cells() if not cell.is_spanned)

def has_merged_cells(table):
    for cell in table.iter_cells():
        if cell.is_merge_origin:
            return True
    return False