from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = ChartData()
chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
chart_data.add_series('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

prs.save('add_chart07.pptx')