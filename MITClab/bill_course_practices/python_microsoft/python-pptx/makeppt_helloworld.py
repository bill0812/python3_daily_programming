import os
from pptx import Presentation

#  創立投影片物件
prs = Presentation()

# 選擇投影片（空白頁、首頁...等等）
title_slide_layout = prs.slide_layouts[0]

# 將該頁加進剛剛創立的投影片物件
slide = prs.slides.add_slide(title_slide_layout)

# 在該頁創立要顯示的內容，包括標題、副標題
title = slide.shapes.title
subtitle = slide.placeholders[1]

# 設定內容
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

# 儲存 PPT
prs.save(os.getcwd() + '/hello_world.pptx')
