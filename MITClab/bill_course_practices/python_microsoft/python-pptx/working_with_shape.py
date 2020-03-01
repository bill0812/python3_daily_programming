from pptx.enum.shapes import MSO_SHAPE
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes
left = top = width = height = Inches(1.0)
shape = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)

from pptx.util import Inches, Pt

length = Inches(1)
print(length)
print(length.inches)
print(length.cm)
print(length.pt)

length = Pt(72)
print(length)

left = top = width = height = Inches(1.0)
shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
print(shape.left, shape.top, shape.width, shape.height)
print(shape.left.inches)

shape.left = Inches(2.0)
print(shape.left.inches)

fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)
fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
fill.fore_color.brightness = -0.25

This sets the shape fill to transparent
shape.fill.background()

line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
line.color.brightness = 0.5  # 50% lighter
line.width = Pt(2.5)

line.color.theme_color = MSO_THEME_COLOR.ACCENT_6

line.fill.solid()
line.fill.fore_color
line.fill.background()

print(line.width,line.width.pt)

line.width = Pt(2.0)
print(line.width.pt)

callout_sp = shapes.add_shape(
    MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR, left, top, width, height
)

# get the callout line coming out of the right place
adjs = callout_sp.adjustments
adjs[0] = 0.5   # vert pos of junction in margin line, 0 is top
adjs[1] = 0.0   # horz pos of margin ln wrt shape width, 0 is left side
adjs[2] = 0.5   # vert pos of elbow wrt margin line, 0 is top
adjs[3] = -0.1  # horz pos of elbow wrt shape width, 0 is margin line
adjs[4] = 3.0   # vert pos of line end wrt shape height, 0 is top
a5 = adjs[3] - (adjs[4] - adjs[0]) * height/width
adjs[5] = a5    # horz pos of elbow wrt shape width, 0 is margin line

# rotate 45 degrees counter-clockwise
callout_sp.rotation = -45.0